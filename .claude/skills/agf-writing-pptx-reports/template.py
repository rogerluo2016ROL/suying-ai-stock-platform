"""
writing-pptx-reports skill — 可复用 helper 全套 + 5 页样例 deck

直接运行（生成 sample-deck.pptx + 自动转 PDF + PNG 预览）：
    python3 template.py

依赖：
    pip3 install --user --break-system-packages python-pptx
    # 可选（验证渲染）：brew install --cask libreoffice && brew install poppler

本文件 = 一份"复制改"型模板，结构对应 SKILL.md 的 12 个关键技巧。
真实项目里可以拷整个文件做起点（如 build-procedure-ppt.py 的祖本）。

设计原则：
- 单一品牌色 + 灰阶（≤ 12 色变量）
- 中文字体 lxml 写 <a:ea>（跨平台不 fallback）
- textbox margin 归零 + word_wrap 显式
- 表格关 firstRow/bandRow + 手动斑马纹
- 卡片化 + 现代 ▎ bullet
- 章节扉页与内容页 layout 分离
"""

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


# ============================================================================
# 1. 设计 token（直接改这里换主题，全 deck 联动）
# ============================================================================

# 字体（默认 Heiti SC — LibreOffice 跨平台稳定渲染）
# ⚠️ 不要默认用 PingFang SC：它在 macOS 是 /System/Library/AssetsV2/... 非标准路径，
#    LibreOffice 找不到 → fallback 到 cursive 花体；前任 build 已实测踩坑。
#    Heiti SC 位于 /System/Library/Fonts/STHeiti Medium.ttc（标准路径），LibreOffice 必识别。
# Windows 端可改 "Microsoft YaHei"；Linux 端可改 "Source Han Sans CN" / "Noto Sans CJK SC"。
FONT_CN = "Heiti SC"
FONT_EN = "Helvetica Neue"
FONT_NUM = "Helvetica Neue"  # 装饰大数字

# 品牌色（这里是 GAC 广汽红样例，换公司就改这 7 行）
PRIMARY_DEEP = RGBColor(0x8B, 0x1F, 0x24)  # 主调（章节扉页 / 表头）
PRIMARY_DARK = RGBColor(0x5E, 0x0E, 0x14)  # 大色块
PRIMARY_TINT = RGBColor(0xFB, 0xE5, 0xE7)  # 装饰底
ACCENT_BRAND = RGBColor(0xEC, 0x0A, 0x1E)  # 强调 / 警示
ACCENT_PALE = RGBColor(0xFE, 0xE1, 0xE4)   # 警示底
GOLD = RGBColor(0xC8, 0xA1, 0x4E)          # 少量装饰短线

# 灰阶（5 档 + 白）
GRAY_900 = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_700 = RGBColor(0x4A, 0x4A, 0x4A)
GRAY_500 = RGBColor(0x8C, 0x8C, 0x8C)
GRAY_300 = RGBColor(0xD9, 0xD9, 0xD9)
GRAY_50 = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# 16:9 页面（13.333 × 7.5 in）
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
LEFT_MARGIN = Inches(0.55)
RIGHT_MARGIN = Inches(0.55)
HEADER_BOTTOM = Inches(1.4)
FOOTER_TOP = Inches(7.0)
CONTENT_W = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN
CONTENT_H = FOOTER_TOP - HEADER_BOTTOM


# ============================================================================
# 2. 字体工具（最致命的坑：必须 lxml 写 <a:ea>）
# ============================================================================

def set_font(run, *, name=FONT_CN, size=14, bold=False, italic=False, color=GRAY_900):
    """
    设置 run 的字体属性 — 必须用 lxml 显式写 <a:ea> + <a:cs>，
    否则中文会 fallback 到丑字体。

    python-pptx 默认 run.font.name 只写 <a:latin>，对中文无效。

    适用：你 `slide.shapes.add_textbox(...)` 自己加的 textbox 的 run。
    Placeholder（layout 自带的标题 / 副标题 / 内容区）请用 `_fix_ph_font(ph, ...)`。
    """
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        elem = rPr.find(qn(tag))
        if elem is None:
            elem = etree.SubElement(rPr, qn(tag))
        elem.set("typeface", name)


def _fix_ph_font(ph, *, name=FONT_CN, size_pt=None, bold=None, italic=None, color=None):
    """
    修复 placeholder 的字体 — 写整个 text_frame 里所有 run 的 <a:ea> + <a:cs>。

    为什么需要：layout / master 的 placeholder 默认继承 master 的 fontScheme，
    很多 iSlide 模板的 master ea 字体是「微软雅黑」—— Windows 字体，
    macOS / LibreOffice 找不到 → 渲染成 cursive 花体。

    set_font(run, ...) 在 slide 级改 run，但 placeholder 的"默认 run"是从 master 继承的，
    必须用本函数直接遍历 placeholder.text_frame 的所有 paragraph.run 写 lxml 节点。

    用法：
        ph.text = "我的标题"
        _fix_ph_font(ph, size_pt=28, bold=True, color=NAVY)

    可选参数全部默认 None：只想修字体不改大小 / 颜色时省略对应参数即可。
    """
    if not ph.has_text_frame:
        return
    tf = ph.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            rPr = run._r.get_or_add_rPr()
            for tag in ("a:latin", "a:ea", "a:cs"):
                elem = rPr.find(qn(tag))
                if elem is None:
                    elem = etree.SubElement(rPr, qn(tag))
                elem.set("typeface", name)
            if size_pt is not None:
                run.font.size = Pt(size_pt)
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic
            if color is not None:
                run.font.color.rgb = color


# ============================================================================
# 2.5 模板工具（基于已有 .pptx 模板生成时使用）
# ============================================================================

def clear_template_slides(prs):
    """
    清空模板自带的所有样例 slide，保留 layout / master / theme。

    为什么需要：仓库 `template/*.pptx` 这类设计模板自带 30+ 张样例 slide
    （含工具说明页 / iSlide 插件指南页等），生成时必须先清空避免污染输出。

    ⚠️ 修复说明：只从 XML sldIdLst 删元素是不够的，OPC relationship
    仍留在 prs.part.rels 里。python-pptx 保存时会遍历所有 rels 写入 ZIP，
    导致旧 slide XML 与新 slide XML 重名 → 40 条 duplicate ZIP entry →
    LibreOffice「source file could not be loaded」。
    正确做法：先 drop_rel 断开 OPC 链接，再清 sldIdLst。

    用法：
        prs = Presentation("template/Template.pptx")
        clear_template_slides(prs)
        # 此刻 prs 是空壳但 layout / master / theme 全保留
        s = prs.slides.add_slide(prs.slide_layouts[0])  # 用第一个 layout 新建 slide
    """
    from pptx.oxml.ns import qn as _qn
    _R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    _SLIDE_RELTYPE = _R_NS + "/slide"

    # 1. 从 XML 直接拿 sldIdLst（不走 prs.slides accessor，避免 rename_slide_parts
    #    在 rels 被清空后找不到 rId 而抛 KeyError）
    sldIdLst = prs.part._element.find(_qn("p:sldIdLst"))
    if sldIdLst is None:
        return

    # 2. 收集每个 sldId 的 r:id，然后清空 XML 列表
    slide_rids = [
        elem.get(f"{{{_R_NS}}}id")
        for elem in list(sldIdLst)
    ]
    for elem in list(sldIdLst):
        sldIdLst.remove(elem)

    # 3. drop_rel 断开 OPC 链接（PackageWriter 遍历 rels 写 ZIP 时不再包含旧 slide Part）
    existing_rids = {rId for rId in prs.part.rels}
    for rId in slide_rids:
        if rId and rId in existing_rids:
            prs.part.drop_rel(rId)


def dump_template_structure(pptx_path):
    """
    打印模板结构概览（layout 清单 + slide → layout 映射）——分析模板时用。

    与 SKILL.md 「Step 1 — 模板分析三件套」配合使用。
    实际写 PPT 时不调用本函数。

    用法：
        python3 -c "from template import dump_template_structure; \\
                    dump_template_structure('template/Template.pptx')"
    """
    prs = Presentation(pptx_path)
    print(f"# Slide size: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inches")
    print(f"# Slides: {len(prs.slides)}  Layouts: {len(prs.slide_layouts)}\n")
    print("## Layouts")
    for i, layout in enumerate(prs.slide_layouts):
        phs = [(ph.placeholder_format.idx, ph.placeholder_format.type) for ph in layout.placeholders]
        print(f"  layout[{i}] '{layout.name}' — {len(phs)} ph, {len(layout.shapes)} sh — {phs}")
    print("\n## Slide → Layout mapping")
    for i, slide in enumerate(prs.slides, 1):
        print(f"  slide{i:2d} → '{slide.slide_layout.name}'")


# ============================================================================
# 3. 形状 helper（真无填充 / 真无边框）
# ============================================================================

def rect(s, x, y, w, h, fill=None, line=None, shape_type=MSO_SHAPE.RECTANGLE):
    """
    画矩形 / 圆角 / 任意 MSO_SHAPE。
    fill=None → 真透明（shape.fill = None 是错的，那是"默认")
    line=None → 真无边框
    line=(color, width) → 自定义边框
    """
    shp = s.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    elif isinstance(line, tuple):
        shp.line.color.rgb = line[0]
        shp.line.width = line[1]
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    return shp


def textbox(s, x, y, w, h, text="", *, size=14, bold=False, italic=False,
            color=GRAY_900, font=FONT_CN, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, line_spacing=1.3, word_wrap=True):
    """
    画文字框 — 所有 margin 归零、line_spacing 显式、word_wrap 显式。
    大字号装饰文本要传 word_wrap=False 防换行。
    """
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    set_font(r, name=font, size=size, bold=bold, italic=italic, color=color)
    return box


# ============================================================================
# 4. 卡片化（圆角矩形 + 左 accent 色条）
# ============================================================================

def card(s, x, y, w, h, *, fill=WHITE, border=GRAY_300, accent=None,
         accent_w=Emu(36000)):
    """
    卡片 = 圆角矩形 + 0.75pt 边框 + 可选左色条。
    accent 给一个 RGBColor，画 4pt 左色条作为视觉锚点。
    """
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.adjustments[0] = 0.05  # 圆角调小（默认大圆角丑）
    if accent:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, accent_w, h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
    return shape


# ============================================================================
# 5. 现代 bullet（▎ 短竖线 + 中文 1.45 行高）
# ============================================================================

def bullets(s, x, y, w, h, items, *, size=14, accent=PRIMARY_DEEP, body=GRAY_900):
    """
    现代制度 PPT 用 ▎ 比 • 更现代。
    items: list[str]
    """
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.45
        r1 = p.add_run()
        r1.text = "▎ "
        set_font(r1, size=size, color=accent, bold=True)
        r2 = p.add_run()
        r2.text = item
        set_font(r2, size=size, color=body)
    return box


# ============================================================================
# 6. 装饰大数字（防换行 + 极淡色）
# ============================================================================

def page_decoration(s, num, *, tint=PRIMARY_TINT):
    """
    右上 140pt 大数字作为页码视觉锚点。
    word_wrap=False 防 "01" → "0\\n1"
    """
    box = s.shapes.add_textbox(Inches(8.8), Inches(0.25), Inches(4.4), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = False  # 关键
    tf.margin_left = tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = str(num)
    set_font(r, name=FONT_NUM, size=140, bold=True, color=tint)
    return box


# ============================================================================
# 7. 现代表格（关 banding + 手工斑马纹 + 显式行高）
# ============================================================================

def table_modern(s, x, y, w, h, headers, rows, *,
                 row_height=Inches(0.4), header_color=PRIMARY_DEEP,
                 alt_row_color=GRAY_50, header_font_color=WHITE,
                 body_font_color=GRAY_900, font_size_header=12,
                 font_size_body=11):
    """
    现代风表格：
    - 表头深色填充
    - 0 内边框（仅顶/底 hairline）
    - 手工斑马纹（关 firstRow/bandRow 防 banding 默认）
    - 显式行高（防 LibreOffice 渲染失控）
    """
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = s.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table

    # 行高（防 LibreOffice 渲染失控）
    for row in tbl.rows:
        row.height = row_height

    # 关 python-pptx 默认 banding
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        tblPr.set("firstRow", "0")
        tblPr.set("bandRow", "0")

    # 表头
    for j, h_text in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.margin_left = tf.margin_right = Emu(45720)
        tf.margin_top = tf.margin_bottom = Emu(36000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h_text
        set_font(r, size=font_size_header, bold=True, color=header_font_color)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 数据行（手工斑马纹）
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i + 1, j)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_row_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Emu(45720)
            tf.margin_top = tf.margin_bottom = Emu(36000)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(cell_text)
            set_font(r, size=font_size_body, color=body_font_color)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return tbl


# ============================================================================
# 8. 页脚 + 页码（每页统一）
# ============================================================================

def footer(s, *, page_num, total, deck_title="样例 deck"):
    """页脚分隔线 + 左 deck 标题 + 右页码 N / TOTAL。"""
    # 分隔线
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              LEFT_MARGIN, FOOTER_TOP, CONTENT_W, Emu(6350))
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY_300
    line.line.fill.background()
    # 左
    textbox(s, LEFT_MARGIN, Inches(7.1), Inches(7.0), Inches(0.35),
            deck_title, size=9, color=GRAY_500)
    # 右
    textbox(s, Inches(11.3), Inches(7.1), Inches(1.5), Inches(0.35),
            f"{page_num} / {total}", size=9, color=GRAY_500,
            align=PP_ALIGN.RIGHT)


def header_line(s):
    """顶部 6pt 极细线分隔（替代厚色带）。"""
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              LEFT_MARGIN, Inches(0.3), CONTENT_W, Emu(76200))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_DEEP
    line.line.fill.background()


# ============================================================================
# 9. 章节扉页 layout（与内容页不同！第 2 反模式）
# ============================================================================

def section_cover(prs, *, section_num, title, subtitle):
    """章节扉页：大色块 + 巨大数字 + 标题 + 副标题。"""
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)
    # 整页淡色底
    rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=PRIMARY_TINT)
    # 左侧大色块
    rect(s, Inches(0.55), Inches(1.5), Inches(2.8), Inches(4.5),
         fill=PRIMARY_DEEP)
    # 巨大数字
    textbox(s, Inches(0.55), Inches(1.5), Inches(2.8), Inches(4.5),
            f"{section_num:02d}", size=180, bold=True, color=WHITE,
            font=FONT_NUM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            word_wrap=False)
    # 标题（右侧）
    textbox(s, Inches(4.0), Inches(3.0), Inches(9.0), Inches(1.0),
            title, size=40, bold=True, color=PRIMARY_DEEP)
    # 短装饰线
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(4.0), Inches(4.0), Inches(1.0), Emu(38100))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    # 副标题
    textbox(s, Inches(4.0), Inches(4.3), Inches(9.0), Inches(1.0),
            subtitle, size=16, color=GRAY_700)
    return s


# ============================================================================
# 10. 内容页 layout（标题 + 装饰大数字 + 内容区）
# ============================================================================

def content_page(prs, *, page_num, title):
    """内容页 layout：顶部细线 + 标题 + 右上装饰数字 + 内容区域 + 页脚。"""
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)
    header_line(s)
    # 标题
    textbox(s, LEFT_MARGIN, Inches(0.55), Inches(8.0), Inches(0.7),
            title, size=24, bold=True, color=PRIMARY_DEEP)
    # 装饰大数字
    page_decoration(s, f"{page_num:02d}", tint=PRIMARY_TINT)
    return s


# ============================================================================
# 11. 警示框（⚠ / ⛔ / 🔒 类，制度 PPT 唯一允许的 emoji）
# ============================================================================

def callout(s, x, y, w, h, text, *, kind="warn"):
    """
    警示框 — kind 决定色调：warn 黄、danger 红、info 蓝。
    制度 PPT 仅允许这种警示性 emoji，不允许 🎉 ✅ 🚀。
    """
    palette = {
        "warn": (RGBColor(0xFF, 0xF7, 0xE0), RGBColor(0xBF, 0x8F, 0x00), "⚠"),
        "danger": (ACCENT_PALE, ACCENT_BRAND, "⛔"),
        "info": (RGBColor(0xE7, 0xF1, 0xFF), RGBColor(0x1F, 0x4E, 0x79), "🔒"),
    }
    bg, fg, icon = palette[kind]
    rect(s, x, y, w, h, fill=bg, line=(fg, Pt(0.75)))
    # 图标
    textbox(s, x, y, Inches(0.6), h, icon, size=20,
            color=fg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 正文
    textbox(s, x + Inches(0.6), y, w - Inches(0.6), h, text,
            size=12, color=GRAY_900, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)


# ============================================================================
# 12. 主函数：生成 5 页样例 deck（封面 / 章节扉页 / bullet / 表格 / 卡片）
# ============================================================================

def build_sample_deck(out_path="sample-deck.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 5

    # P1 — 封面
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=PRIMARY_DEEP)
    textbox(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.4),
            "writing-pptx-reports", size=54, bold=True, color=WHITE,
            font=FONT_EN)
    textbox(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.6),
            "现代化中文制度 PPT 实战模板", size=22, color=PRIMARY_TINT)
    # 装饰线
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.8), Inches(4.7), Inches(2.0), Emu(38100))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    textbox(s, Inches(0.8), Inches(4.95), Inches(11.5), Inches(0.5),
            "v1.0 · 2026-05 · GAC 红主调样例", size=12, color=PRIMARY_TINT)

    # P2 — 章节扉页
    section_cover(prs, section_num=1, title="设计原则",
                  subtitle="单一品牌色 + 卡片化 + 现代 bullet")

    # P3 — bullet 内容页
    s = content_page(prs, page_num=2, title="核心 7 条")
    bullets(s, LEFT_MARGIN, Inches(1.5), CONTENT_W, CONTENT_H,
            ["中文字体必须 lxml 写 <a:ea>，否则跨平台 fallback",
             "textbox margin 全归零，否则文字神秘偏右",
             "大字号 textbox 必须 word_wrap=False，否则 \"01\" 换行",
             "表格关 firstRow/bandRow，手工斑马纹",
             "shape 真无填充要 fill.background()，None 是\"默认\"",
             "▎ 比 • 更现代，制度 PPT 标配",
             "章节扉页与内容页 layout 必须不同"], size=16)
    footer(s, page_num=2, total=TOTAL)

    # P4 — 表格内容页
    s = content_page(prs, page_num=3, title="字号层级（16:9 = 13.333 × 7.5 in）")
    table_modern(s, LEFT_MARGIN, Inches(1.6), CONTENT_W, Inches(4.5),
                 ["用途", "字号", "加粗", "示例场景"],
                 [["封面主标题", "44-54pt", "bold", "deck 第一页"],
                  ["章节扉页大标题", "36-40pt", "bold", "分章页"],
                  ["内容页 H2", "20-28pt", "bold", "页面标题"],
                  ["内容页 H3", "14-18pt", "bold", "小节标题"],
                  ["正文 bullet", "11.5-14pt", "normal", "正文清单"],
                  ["表格 body", "10.5-12pt", "normal", "表格单元格"],
                  ["页脚 / caption", "8.5-10pt", "normal", "页脚标注"],
                  ["装饰大数字", "120-150pt", "bold（淡色）", "右上锚点"]])
    footer(s, page_num=3, total=TOTAL)

    # P5 — 卡片化 + 警示框
    s = content_page(prs, page_num=4, title="卡片化 + 警示框")
    points = [
        ("视觉锚点", "左 4pt 色条 = 视觉锚点，引导眼睛"),
        ("圆角调小", "adjustments[0]=0.05；默认大圆角丑"),
        ("0.75pt 边框", "GRAY_300，比 1pt 内敛"),
        ("Padding", "卡片内文字距边界 0.15 inch 起步"),
    ]
    for i, (k, v) in enumerate(points):
        y = Inches(1.7 + i * 0.85)
        card(s, LEFT_MARGIN, y, Inches(8.0), Inches(0.7),
             fill=GRAY_50, border=GRAY_300, accent=PRIMARY_DEEP)
        textbox(s, Inches(0.8), y, Inches(2.5), Inches(0.7),
                k, size=14, bold=True, color=PRIMARY_DEEP,
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, Inches(3.5), y, Inches(4.4), Inches(0.7),
                v, size=12, color=GRAY_700, anchor=MSO_ANCHOR.MIDDLE)
    # 右侧 3 个警示框
    callout(s, Inches(9.0), Inches(1.7), Inches(3.8), Inches(0.9),
            "PrinceXML 商业 / 免费版有水印", kind="danger")
    callout(s, Inches(9.0), Inches(2.85), Inches(3.8), Inches(0.9),
            "默认 font.name 对中文无效", kind="warn")
    callout(s, Inches(9.0), Inches(4.0), Inches(3.8), Inches(0.9),
            "只用 python-pptx 读回验证不够", kind="info")
    footer(s, page_num=4, total=TOTAL)

    # 保存
    out = Path(out_path).resolve()
    prs.save(str(out))
    print(f"✅ 已生成 {out}")
    print(f"\n推荐立即验证（实际渲染往往与 python-pptx 内存视图不同）：")
    print(f"  cd /tmp && rm -rf preview && mkdir preview && cd preview")
    print(f"  soffice --headless --convert-to pdf {out}")
    print(f"  pdftocairo -png -r 100 {out.stem}.pdf p")
    print(f"  open p-*.png")
    return out


if __name__ == "__main__":
    build_sample_deck()
